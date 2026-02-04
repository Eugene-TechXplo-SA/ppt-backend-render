from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import tempfile
import os
import zipfile
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import Image
import re
import logging
import io
from urllib.parse import urlparse
from typing import Dict, List, Optional, Tuple
import unicodedata
from difflib import SequenceMatcher

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

def normalize_name(name: str) -> str:
    """Normalize a name for flexible matching: lowercase, remove special chars and ALL spaces."""
    if not name:
        return ""
    normalized = unicodedata.normalize('NFKD', name)
    normalized = normalized.encode('ASCII', 'ignore').decode('ASCII')
    normalized = re.sub(r'[^a-zA-Z0-9\s]', '', normalized)
    normalized = re.sub(r'\s+', '', normalized)
    return normalized.lower()

def find_best_column_match(field_name: str, columns: List[str], threshold: float = 0.85) -> Optional[str]:
    """
    Find the best matching column for a field name.
    First tries exact normalized match, then falls back to fuzzy matching.

    Args:
        field_name: The field name from the placeholder
        columns: List of available column names
        threshold: Minimum similarity ratio for fuzzy matching (0-1)

    Returns:
        The best matching column name, or None if no good match found
    """
    field_normalized = normalize_name(field_name)

    # Try exact normalized match first
    for col in columns:
        if normalize_name(col) == field_normalized:
            logger.debug(f"Exact match found: '{field_name}' -> '{col}'")
            return col

    # Fall back to fuzzy matching
    best_match = None
    best_ratio = 0.0

    for col in columns:
        col_normalized = normalize_name(col)
        ratio = SequenceMatcher(None, field_normalized, col_normalized).ratio()

        if ratio > best_ratio and ratio >= threshold:
            best_ratio = ratio
            best_match = col

    if best_match:
        logger.info(f"Fuzzy match found: '{field_name}' -> '{best_match}' (similarity: {best_ratio:.2%})")
    else:
        logger.debug(f"No match found for: '{field_name}'")

    return best_match

def calculate_fit_dimensions(img_path: str, shape_width: int, shape_height: int, shape_left: int, shape_top: int) -> Tuple[int, int, int, int]:
    """
    Calculate dimensions and position to fit an image within a shape while maintaining aspect ratio.
    Returns: (left, top, width, height) for the fitted image.
    """
    try:
        with Image.open(img_path) as img:
            img_width, img_height = img.size

        # Calculate aspect ratios
        img_aspect = img_width / img_height
        shape_aspect = shape_width / shape_height

        # Determine dimensions based on aspect ratio comparison
        if img_aspect > shape_aspect:
            # Image is wider - constrain by width
            new_width = shape_width
            new_height = int(shape_width / img_aspect)
        else:
            # Image is taller or equal - constrain by height
            new_height = shape_height
            new_width = int(shape_height * img_aspect)

        # Center the image within the shape
        new_left = shape_left + (shape_width - new_width) // 2
        new_top = shape_top + (shape_height - new_height) // 2

        return new_left, new_top, new_width, new_height

    except Exception as e:
        logger.error(f"Error calculating fit dimensions: {str(e)}")
        # Fallback to original dimensions if calculation fails
        return shape_left, shape_top, shape_width, shape_height

def find_folder_for_site(site_identifier: str, images_dir: str) -> Optional[str]:
    """Find a folder matching the site identifier using flexible matching."""
    if not site_identifier or not os.path.exists(images_dir):
        return None

    normalized_site = normalize_name(site_identifier)

    for item in os.listdir(images_dir):
        item_path = os.path.join(images_dir, item)
        if os.path.isdir(item_path):
            normalized_folder = normalize_name(item)
            if normalized_folder == normalized_site:
                logger.info(f"Matched site '{site_identifier}' to folder '{item}'")
                return item_path

    return None

def get_first_image_from_folder(folder_path: str) -> Optional[str]:
    """Get the first image file from a folder (alphabetically)."""
    if not os.path.exists(folder_path):
        return None

    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}
    images = []

    for file in os.listdir(folder_path):
        file_lower = file.lower()
        if any(file_lower.endswith(ext) for ext in image_extensions):
            images.append(file)

    if images:
        images.sort()
        return os.path.join(folder_path, images[0])

    return None

def detect_image_type_from_column(column_name: str) -> str:
    """
    Detect the image type based on column name.
    Returns: 'map_screenshot', 'site_image', or 'unknown'
    """
    col_lower = column_name.lower()

    # Check for map/screenshot keywords
    if any(keyword in col_lower for keyword in ['map', 'screenshot', 'location']):
        return 'map_screenshot'

    # Check for site image keywords
    if any(keyword in col_lower for keyword in ['site', 'photo', 'picture', 'image', 'pic']):
        return 'site_image'

    return 'unknown'

def find_subfolder(site_folder: str, subfolder_type: str) -> Optional[str]:
    """
    Find a subfolder within a site folder based on type.
    Supports flexible matching for subfolder names.

    Args:
        site_folder: Path to the site folder
        subfolder_type: Either 'site_image' or 'map_screenshot'

    Returns:
        Path to the subfolder if found, None otherwise
    """
    if not os.path.exists(site_folder) or not os.path.isdir(site_folder):
        return None

    # Define possible subfolder name variations
    subfolder_patterns = {
        'site_image': ['site image', 'siteimage', 'site_image', 'site-image', 'site'],
        'map_screenshot': ['map screenshot', 'mapscreenshot', 'map_screenshot', 'map-screenshot', 'map']
    }

    patterns = subfolder_patterns.get(subfolder_type, [])

    for item in os.listdir(site_folder):
        item_path = os.path.join(site_folder, item)
        if os.path.isdir(item_path):
            normalized_item = normalize_name(item)
            for pattern in patterns:
                if normalized_item == normalize_name(pattern):
                    logger.info(f"Found subfolder '{item}' matching type '{subfolder_type}'")
                    return item_path

    return None

def get_image_from_subfolder(site_folder: str, image_type: str) -> Optional[str]:
    """
    Get the first image from the appropriate subfolder within a site folder.

    Args:
        site_folder: Path to the site folder
        image_type: Either 'map_screenshot' or 'site_image'

    Returns:
        Path to the image if found, None otherwise
    """
    if not site_folder or not image_type:
        return None

    subfolder = find_subfolder(site_folder, image_type)
    if subfolder:
        image_path = get_first_image_from_folder(subfolder)
        if image_path:
            logger.info(f"Found image in subfolder '{subfolder}': {image_path}")
            return image_path
        else:
            logger.warning(f"Subfolder '{subfolder}' exists but contains no images")
    else:
        logger.debug(f"No subfolder found for type '{image_type}' in '{site_folder}'")

    return None

def find_image_path_enhanced(value: str, images_dir: str, site_identifier: Optional[str] = None, image_type: str = 'unknown') -> Optional[str]:
    """
    Enhanced image finding with subfolder-based support.

    Strategy:
    1. If site_identifier and image_type provided, look for image in appropriate subfolder
    2. Fall back to legacy folder structure (all images in site folder)
    3. Fall back to exact filename matching (backward compatibility)
    4. Fall back to case-insensitive filename search

    Args:
        value: The value from the Excel cell
        images_dir: Root directory containing images
        site_identifier: Name/ID of the site
        image_type: Type of image ('site_image', 'map_screenshot', or 'unknown')
    """
    try:
        if site_identifier:
            folder_path = find_folder_for_site(site_identifier, images_dir)
            if folder_path:
                # Try new subfolder structure first
                if image_type != 'unknown':
                    image_path = get_image_from_subfolder(folder_path, image_type)
                    if image_path:
                        logger.info(f"Found image in subfolder (type: {image_type}): {image_path}")
                        return image_path
                    else:
                        logger.info(f"No image found in subfolder for type '{image_type}', trying legacy structure")

                # Fall back to legacy structure (images directly in site folder)
                image_path = get_first_image_from_folder(folder_path)
                if image_path:
                    logger.info(f"Found image in legacy folder structure: {image_path}")
                    return image_path

        clean_value = value.replace("images/", "", 1) if value.startswith("images/") else value
        candidate = os.path.join(images_dir, clean_value)
        if os.path.exists(candidate):
            logger.info(f"Found image by exact path: {candidate}")
            return candidate

        for root, dirs, files in os.walk(images_dir):
            for file in files:
                if file.lower() == clean_value.lower():
                    found_path = os.path.join(root, file)
                    logger.info(f"Found image by case-insensitive search: {found_path}")
                    return found_path

        logger.warning(f"No image found for value: {value}")
        return None
    except Exception as e:
        logger.error(f"Error in find_image_path_enhanced: {str(e)}")
        return None

def validate_zip_structure(zip_path: str, df: pd.DataFrame, site_column: Optional[str] = None) -> Dict:
    """Validate ZIP structure and return detailed report with subfolder support."""
    validation_result = {
        "valid": True,
        "errors": [],
        "warnings": [],
        "structure": {
            "has_folders": False,
            "folders": [],
            "loose_files": [],
            "matched_sites": [],
            "unmatched_sites": [],
            "extra_folders": [],
            "subfolder_structure": {}
        }
    }

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            extract_dir = os.path.join(tmpdir, "extracted")
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)

            items = os.listdir(extract_dir)
            folders = [item for item in items if os.path.isdir(os.path.join(extract_dir, item))]
            files = [item for item in items if os.path.isfile(os.path.join(extract_dir, item))]

            validation_result["structure"]["has_folders"] = len(folders) > 0
            validation_result["structure"]["folders"] = folders
            validation_result["structure"]["loose_files"] = files

            if site_column and site_column in df.columns:
                sites = df[site_column].dropna().unique().tolist()
                sites = [str(s).strip() for s in sites]

                for site in sites:
                    folder_path = find_folder_for_site(site, extract_dir)
                    if folder_path:
                        # Check for subfolder structure
                        site_image_subfolder = find_subfolder(folder_path, 'site_image')
                        map_screenshot_subfolder = find_subfolder(folder_path, 'map_screenshot')

                        has_subfolders = site_image_subfolder or map_screenshot_subfolder
                        site_image_count = 0
                        map_screenshot_count = 0

                        if site_image_subfolder:
                            site_images = [f for f in os.listdir(site_image_subfolder)
                                          if os.path.isfile(os.path.join(site_image_subfolder, f)) and
                                          any(f.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp'])]
                            site_image_count = len(site_images)

                        if map_screenshot_subfolder:
                            map_images = [f for f in os.listdir(map_screenshot_subfolder)
                                         if os.path.isfile(os.path.join(map_screenshot_subfolder, f)) and
                                         any(f.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp'])]
                            map_screenshot_count = len(map_images)

                        # Check for legacy structure (images directly in folder)
                        direct_images = [f for f in os.listdir(folder_path)
                                        if os.path.isfile(os.path.join(folder_path, f)) and
                                        any(f.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp'])]

                        validation_result["structure"]["subfolder_structure"][site] = {
                            "has_subfolders": has_subfolders,
                            "site_image_subfolder": site_image_subfolder is not None,
                            "site_image_count": site_image_count,
                            "map_screenshot_subfolder": map_screenshot_subfolder is not None,
                            "map_screenshot_count": map_screenshot_count,
                            "direct_images": len(direct_images),
                            "structure_type": "subfolder" if has_subfolders else "legacy"
                        }

                        if has_subfolders:
                            if site_image_count > 0 or map_screenshot_count > 0:
                                validation_result["structure"]["matched_sites"].append(site)
                                if site_image_count > 0 and map_screenshot_count > 0:
                                    logger.info(f"✓ Site '{site}': Using subfolder structure with {site_image_count} site image(s) and {map_screenshot_count} map screenshot(s)")
                                elif site_image_count > 0:
                                    validation_result["warnings"].append(f"Site '{site}' has Site Image subfolder but missing Map Screenshot subfolder")
                                elif map_screenshot_count > 0:
                                    validation_result["warnings"].append(f"Site '{site}' has Map Screenshot subfolder but missing Site Image subfolder")
                            else:
                                validation_result["structure"]["unmatched_sites"].append(site)
                                validation_result["warnings"].append(f"Site '{site}' has subfolders but they contain no images")
                        else:
                            if len(direct_images) > 0:
                                validation_result["structure"]["matched_sites"].append(site)
                                validation_result["warnings"].append(f"Site '{site}' is using legacy structure ({len(direct_images)} image(s) directly in folder)")
                            else:
                                validation_result["structure"]["unmatched_sites"].append(site)
                                validation_result["warnings"].append(f"Folder found for '{site}' but no images inside")
                    else:
                        validation_result["structure"]["unmatched_sites"].append(site)
                        validation_result["warnings"].append(f"No folder found for site: {site}")

                normalized_sites = {normalize_name(s) for s in sites}
                for folder in folders:
                    if normalize_name(folder) not in normalized_sites:
                        validation_result["structure"]["extra_folders"].append(folder)
                        validation_result["warnings"].append(f"Extra folder not referenced in Excel: {folder}")

            if len(folders) == 0 and len(files) == 0:
                validation_result["valid"] = False
                validation_result["errors"].append("ZIP file is empty")

    except Exception as e:
        validation_result["valid"] = False
        validation_result["errors"].append(f"Error validating ZIP: {str(e)}")

    return validation_result

def get_value_for_field(row, field):
    try:
        if field in row and not pd.isna(row[field]):
            return str(row[field]).strip()
        return ""
    except:
        return ""

def is_image_path(val):
    return val and isinstance(val, str) and bool(re.search(r'\.(jpg|jpeg|png|gif|bmp)$', val.lower()))

def replace_text_in_obj(obj, row, images_dir, site_identifier=None):
    placeholder_pattern = re.compile(r"\{\{\s*(.*?)\s*\}\}")
    try:
        if hasattr(obj, "text_frame") and obj.text_frame is not None:
            for paragraph in obj.text_frame.paragraphs:
                for run in paragraph.runs:
                    matches = placeholder_pattern.findall(run.text)
                    for field_raw in matches:
                        field_clean = re.sub(r'\s+', ' ', field_raw.strip())

                        col = find_best_column_match(field_clean, row.index.tolist())
                        if not col:
                            continue
                        val = get_value_for_field(row, col)

                        if is_image_path(val):
                            continue

                        escaped = re.escape(f"{{{{{field_raw}}}}}")
                        flexible = (
                            escaped
                            .replace('\\{\\{', r'\s*\{\{\s*')
                            .replace('\\}\\}', r'\s*\}\}')
                            .replace('\\-', r'[-–—]')
                            .replace('\\ ', r'\s*')
                        )
                        match = re.search(flexible, run.text, re.UNICODE)
                        if not match:
                            continue
                        placeholder_text = match.group(0)

                        if col.lower().endswith("link"):
                            try:
                                result = urlparse(val)
                                if all([result.scheme, result.netloc]):
                                    run.text = run.text.replace(placeholder_text, val, 1)
                                    run.font.color.rgb = RGBColor(0, 0, 255)
                                    run.font.underline = True
                                    continue
                            except:
                                pass

                        run.text = run.text.replace(placeholder_text, val, 1)

    except Exception as e:
        logger.error(f"Error in replace_text_in_obj: {str(e)}")

def replace_images_on_shape(shape, row, images_dir, site_identifier=None):
    placeholder_pattern = re.compile(r"\{\{\s*(.*?)\s*\}\}")
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            full_text = "".join(r.text or "" for p in shape.text_frame.paragraphs for r in p.runs)
            matches = placeholder_pattern.findall(full_text)
            for field_raw in matches:
                field_clean = re.sub(r'\s+', ' ', field_raw.strip())

                col = find_best_column_match(field_clean, row.index.tolist())
                if not col:
                    logger.debug(f"No column found for placeholder: {field_raw}")
                    continue

                val = get_value_for_field(row, col)

                is_image_column = any(keyword in col.lower() for keyword in ['image', 'screenshot', 'photo', 'picture', 'pic', 'map'])

                if not is_image_column:
                    logger.debug(f"Skipping non-image column: {col}")
                    continue

                # Detect image type from column name
                image_type = detect_image_type_from_column(col)
                logger.info(f"Processing image placeholder '{field_raw}' (column: '{col}', type: '{image_type}') with value: '{val}' and site: '{site_identifier}'")

                img_path = None

                is_placeholder_text = val and '{{' in str(val) and '}}' in str(val)
                is_empty_or_placeholder = not val or is_placeholder_text

                if is_image_path(val) and not is_placeholder_text:
                    img_path = find_image_path_enhanced(val, images_dir, site_identifier, image_type)
                    if img_path:
                        logger.info(f"Found image via file path: {img_path}")
                elif is_empty_or_placeholder and site_identifier:
                    logger.info(f"Cell is empty/placeholder for image column, trying subfolder lookup for site: {site_identifier} (type: {image_type})")
                    folder_path = find_folder_for_site(site_identifier, images_dir)
                    if folder_path:
                        # Try subfolder structure first
                        if image_type != 'unknown':
                            img_path = get_image_from_subfolder(folder_path, image_type)
                            if img_path:
                                logger.info(f"Found image via subfolder (type: {image_type}): {img_path}")
                            else:
                                logger.info(f"No subfolder found for type '{image_type}', trying legacy structure")

                        # Fall back to legacy structure if subfolder approach didn't work
                        if not img_path:
                            img_path = get_first_image_from_folder(folder_path)
                            if img_path:
                                logger.info(f"Found image via legacy folder structure: {img_path}")
                            else:
                                logger.warning(f"Site folder found but no images inside: {folder_path}")
                    else:
                        logger.warning(f"No folder found for site: {site_identifier}")
                else:
                    logger.debug(f"Cell has non-image value: {val}")

                if img_path:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = r.text.replace(f"{{{{{field_raw}}}}}", "", 1)
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height

                    # Calculate fitted dimensions to maintain aspect ratio
                    fitted_left, fitted_top, fitted_width, fitted_height = calculate_fit_dimensions(
                        img_path, width, height, left, top
                    )

                    sp = shape._element
                    sp.getparent().remove(sp)
                    slide = shape.part.slide
                    slide.shapes.add_picture(img_path, fitted_left, fitted_top, width=fitted_width, height=fitted_height)
                    logger.info(f"Successfully inserted image for placeholder: {field_raw} with fit dimensions")
                    return
                else:
                    logger.warning(f"No image found for placeholder: {field_raw}")
    except Exception as e:
        logger.error(f"Error in replace_images_on_shape: {str(e)}")

@app.post("/api/validate")
async def validate(excel: UploadFile = File(...), images: UploadFile = File(None)):
    """Validate uploaded files before processing."""
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            excel_path = os.path.join(tmpdir, excel.filename or "content.xlsx")
            with open(excel_path, "wb") as f:
                f.write(await excel.read())

            df = pd.read_excel(excel_path)
            df.columns = [col.strip() for col in df.columns]

            result = {
                "valid": True,
                "excel": {
                    "rows": len(df),
                    "columns": list(df.columns),
                    "sample_data": df.head(3).to_dict(orient="records")
                },
                "images": None
            }

            if images:
                zip_path = os.path.join(tmpdir, images.filename or "images.zip")
                images_content = await images.read()
                with open(zip_path, "wb") as f:
                    f.write(images_content)

                # Use same logic as generate endpoint to detect site column
                site_column = None
                priority_keywords = [
                    ['site number', 'site_number'],  # Highest priority
                    ['number'],  # Number column
                    ['site name', 'site_name'],  # Site names
                    ['address', 'city', 'town'],  # Location identifiers
                    ['name', 'site', 'id']  # Generic fallback
                ]

                # Try each priority group and find the first column with good uniqueness
                for keyword_group in priority_keywords:
                    for col in df.columns:
                        col_lower = col.lower()
                        for keyword in keyword_group:
                            if keyword in col_lower:
                                unique_values = df[col].dropna().unique()
                                unique_ratio = len(unique_values) / len(df[col].dropna()) if len(df[col].dropna()) > 0 else 0

                                logger.info(f"Validation - evaluating '{col}': unique ratio {unique_ratio:.2f}")

                                # Accept columns with at least 30% uniqueness
                                if unique_ratio > 0.3:
                                    site_column = col
                                    logger.info(f"Validation - detected site column: '{site_column}' (unique ratio: {unique_ratio:.2f})")
                                    break
                        if site_column:
                            break
                    if site_column:
                        break

                if not site_column:
                    for col in df.columns:
                        if 'site' in col.lower() or 'name' in col.lower():
                            site_column = col
                            logger.info(f"Validation - using fallback column: '{site_column}'")
                            break

                validation = validate_zip_structure(zip_path, df, site_column)
                result["images"] = validation
                result["valid"] = validation["valid"]

            return result

        except Exception as e:
            logger.error(f"Error in /api/validate: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Validation Error: {str(e)}")

@app.post("/api/generate")
async def generate(
    excel: UploadFile = File(...),
    ppt: UploadFile = File(...),
    images: UploadFile = File(None),
    site_column: Optional[str] = None
):
    """Generate PowerPoint with enhanced folder-based image support."""
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            excel_path = os.path.join(tmpdir, excel.filename or "content.xlsx")
            ppt_path = os.path.join(tmpdir, ppt.filename or "template.pptx")
            images_dir = os.path.join(tmpdir, "images")

            with open(excel_path, "wb") as f:
                f.write(await excel.read())
            with open(ppt_path, "wb") as f:
                f.write(await ppt.read())

            images_content = None
            if images:
                images_content = await images.read()
                zip_path = os.path.join(tmpdir, images.filename or "images.zip")
                with open(zip_path, "wb") as f:
                    f.write(images_content)
                with zipfile.ZipFile(zip_path, "r") as zip_ref:
                    zip_ref.extractall(images_dir)
                logger.info(f"Extracted images to: {images_dir}")

                folders = [item for item in os.listdir(images_dir) if os.path.isdir(os.path.join(images_dir, item))]
                files = [item for item in os.listdir(images_dir) if os.path.isfile(os.path.join(images_dir, item))]
                logger.info(f"ZIP structure - Folders: {folders}")
                logger.info(f"ZIP structure - Root files: {files}")

            df = pd.read_excel(excel_path)
            df.columns = [col.strip() for col in df.columns]
            logger.info(f"Excel columns: {list(df.columns)}")

            if not site_column:
                # Try to find the best column for site identification
                # Priority: site number (most unique) > site name > address/city > generic
                priority_keywords = [
                    ['site number', 'site_number'],  # Highest priority - most unique identifiers
                    ['number'],  # Number column
                    ['site name', 'site_name'],  # Site names
                    ['address', 'city', 'town'],  # Location-based identifiers
                    ['name', 'site', 'id']  # Generic fallback
                ]

                # Try each priority group and find the first column with good uniqueness
                for keyword_group in priority_keywords:
                    for col in df.columns:
                        col_lower = col.lower()
                        for keyword in keyword_group:
                            if keyword in col_lower:
                                # Check uniqueness
                                unique_values = df[col].dropna().unique()
                                unique_ratio = len(unique_values) / len(df[col].dropna()) if len(df[col].dropna()) > 0 else 0

                                logger.info(f"Evaluating column '{col}' - unique ratio: {unique_ratio:.2f}")

                                # Accept columns with at least 30% uniqueness
                                if unique_ratio > 0.3:
                                    site_column = col
                                    logger.info(f"✓ Auto-detected site column: '{site_column}' (unique ratio: {unique_ratio:.2f})")
                                    break
                        if site_column:
                            break
                    if site_column:
                        break

                if not site_column:
                    # Last resort - use first column with 'site' or 'name'
                    for col in df.columns:
                        if 'site' in col.lower() or 'name' in col.lower():
                            site_column = col
                            logger.warning(f"Using fallback site column: '{site_column}'")
                            break
            else:
                logger.info(f"Using provided site column: '{site_column}'")

            if site_column and site_column in df.columns:
                logger.info(f"Site values in data: {df[site_column].tolist()}")
            else:
                logger.warning("No site column detected or column not found in data")

            prs = Presentation(ppt_path)

            logger.info("Processing slides – Step 1: Images")
            logger.info(f"Site column being used: '{site_column}'")
            logger.info(f"Site values in data: {df[site_column].tolist() if site_column and site_column in df.columns else 'None'}")

            for i, row in df.iterrows():
                if i >= len(prs.slides):
                    break
                slide = prs.slides[i]

                site_id = None
                if site_column and site_column in row.index:
                    site_id = get_value_for_field(row, site_column)
                    logger.info(f"━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
                    logger.info(f"Slide {i}: Processing site '{site_id}' (from column '{site_column}')")

                    # Check if folder exists for this site
                    if images:
                        test_folder = find_folder_for_site(site_id, images_dir)
                        if test_folder:
                            logger.info(f"✓ Found matching folder: {os.path.basename(test_folder)}")
                        else:
                            logger.warning(f"✗ No folder found matching '{site_id}'")
                            # List available folders to help debug
                            available_folders = [f for f in os.listdir(images_dir) if os.path.isdir(os.path.join(images_dir, f))]
                            logger.warning(f"Available folders: {available_folders}")

                for shape in slide.shapes:
                    replace_images_on_shape(shape, row, images_dir if images else tmpdir, site_id)
                    if hasattr(shape, "table"):
                        for row_cells in shape.table.rows:
                            for cell in row_cells.cells:
                                replace_images_on_shape(cell, row, images_dir if images else tmpdir, site_id)

            logger.info("Processing slides – Step 2: Text")
            for i, row in df.iterrows():
                if i >= len(prs.slides):
                    break
                slide = prs.slides[i]

                site_id = None
                if site_column and site_column in row.index:
                    site_id = get_value_for_field(row, site_column)

                for shape in slide.shapes:
                    replace_text_in_obj(shape, row, images_dir if images else tmpdir, site_id)
                    if hasattr(shape, "table"):
                        for row_cells in shape.table.rows:
                            for cell in row_cells.cells:
                                replace_text_in_obj(cell, row, images_dir if images else tmpdir, site_id)

            output_file = os.path.join(tmpdir, "Client_Presentation.pptx")
            prs.save(output_file)

            with open(output_file, "rb") as f:
                file_content = f.read()

            return StreamingResponse(
                io.BytesIO(file_content),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": "attachment; filename=Client_Presentation.pptx"}
            )

        except Exception as e:
            logger.error(f"Error in /api/generate: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Internal Server Error: {str(e)}")

@app.get("/health")
async def health():
    return {"status": "ok"}
